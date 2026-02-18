"""
Tests for memctl.extract — unified text extraction module.

Tests cover:
    - Extension set consistency (BINARY_EXTS ⊂ ALL_INGESTABLE_EXTS)
    - Text file reading (always available)
    - Binary format dispatch and ImportError handling
    - PDF extraction via pdftotext (if available)
    - Office format extraction (if libraries installed)

Author: Olivier Vitrac, PhD, HDR | olivier.vitrac@adservio.fr | Adservio
"""

import os
import shutil
import subprocess
import tempfile

import pytest

from memctl.extract import (
    ALL_INGESTABLE_EXTS,
    BINARY_EXTS,
    read_file_text,
    _extract_binary,
)


# ---------------------------------------------------------------------------
# Extension sets
# ---------------------------------------------------------------------------

class TestExtensionSets:
    """Verify extension classification invariants."""

    def test_binary_exts_subset_of_all(self):
        """Every binary ext must also be in ALL_INGESTABLE_EXTS."""
        assert BINARY_EXTS.issubset(ALL_INGESTABLE_EXTS)

    def test_all_exts_start_with_dot(self):
        for ext in ALL_INGESTABLE_EXTS:
            assert ext.startswith("."), f"Extension missing dot: {ext!r}"

    def test_binary_exts_start_with_dot(self):
        for ext in BINARY_EXTS:
            assert ext.startswith("."), f"Extension missing dot: {ext!r}"

    def test_binary_exts_expected(self):
        """Verify the 7 binary formats are present."""
        expected = {".docx", ".odt", ".pptx", ".odp", ".xlsx", ".ods", ".pdf"}
        assert expected == BINARY_EXTS

    def test_text_exts_not_binary(self):
        """Common text formats must NOT be in BINARY_EXTS."""
        text_exts = {".md", ".txt", ".py", ".json", ".html", ".csv"}
        assert text_exts.isdisjoint(BINARY_EXTS)

    def test_source_code_exts_present(self):
        """Common source code extensions must be ingestable."""
        code_exts = {".py", ".js", ".ts", ".java", ".go", ".rs", ".c", ".cpp", ".sh"}
        assert code_exts.issubset(ALL_INGESTABLE_EXTS)

    def test_office_exts_present(self):
        """All Office/ODF extensions must be ingestable."""
        office = {".docx", ".odt", ".pptx", ".odp", ".xlsx", ".ods"}
        assert office.issubset(ALL_INGESTABLE_EXTS)


# ---------------------------------------------------------------------------
# Text file reading
# ---------------------------------------------------------------------------

class TestReadTextFiles:
    """Test direct text file reading (no external libraries needed)."""

    def test_read_markdown(self, tmp_path):
        f = tmp_path / "test.md"
        f.write_text("# Hello\n\nWorld\n", encoding="utf-8")
        text = read_file_text(str(f))
        assert "# Hello" in text
        assert "World" in text

    def test_read_python(self, tmp_path):
        f = tmp_path / "test.py"
        f.write_text("def foo():\n    return 42\n", encoding="utf-8")
        text = read_file_text(str(f))
        assert "def foo():" in text

    def test_read_json(self, tmp_path):
        f = tmp_path / "test.json"
        f.write_text('{"key": "value"}', encoding="utf-8")
        text = read_file_text(str(f))
        assert '"key"' in text

    def test_read_html(self, tmp_path):
        f = tmp_path / "test.html"
        f.write_text("<html><body>Hello</body></html>", encoding="utf-8")
        text = read_file_text(str(f))
        assert "Hello" in text

    def test_read_csv(self, tmp_path):
        f = tmp_path / "data.csv"
        f.write_text("a,b,c\n1,2,3\n", encoding="utf-8")
        text = read_file_text(str(f))
        assert "a,b,c" in text

    def test_nonexistent_file_raises(self, tmp_path):
        with pytest.raises(FileNotFoundError):
            read_file_text(str(tmp_path / "nope.txt"))

    def test_encoding_errors_replaced(self, tmp_path):
        """Binary garbage in a .txt file should not raise."""
        f = tmp_path / "bad.txt"
        f.write_bytes(b"Hello \xff\xfe World")
        text = read_file_text(str(f))
        assert "Hello" in text
        assert "World" in text

    def test_empty_file(self, tmp_path):
        f = tmp_path / "empty.txt"
        f.write_text("", encoding="utf-8")
        text = read_file_text(str(f))
        assert text == ""


# ---------------------------------------------------------------------------
# Binary dispatch and ImportError
# ---------------------------------------------------------------------------

class TestBinaryDispatch:
    """Test the binary extraction dispatcher."""

    def test_unsupported_ext_raises(self, tmp_path):
        f = tmp_path / "data.xyz"
        f.write_bytes(b"stuff")
        with pytest.raises(ValueError, match="Unsupported binary format"):
            _extract_binary(str(f), ".xyz")

    def test_docx_import_error(self, tmp_path):
        """If python-docx is not installed, ImportError mentions install."""
        try:
            import docx  # noqa: F401
            pytest.skip("python-docx is installed — skip ImportError test")
        except ImportError:
            pass
        f = tmp_path / "test.docx"
        f.write_bytes(b"PK\x03\x04")  # minimal ZIP header
        with pytest.raises(ImportError, match="pip install.*memctl"):
            read_file_text(str(f))

    def test_odt_import_error(self, tmp_path):
        try:
            import odf  # noqa: F401
            pytest.skip("odfpy is installed — skip ImportError test")
        except ImportError:
            pass
        f = tmp_path / "test.odt"
        f.write_bytes(b"PK\x03\x04")
        with pytest.raises(ImportError, match="pip install.*memctl"):
            read_file_text(str(f))

    def test_pptx_import_error(self, tmp_path):
        try:
            import pptx  # noqa: F401
            pytest.skip("python-pptx is installed — skip ImportError test")
        except ImportError:
            pass
        f = tmp_path / "test.pptx"
        f.write_bytes(b"PK\x03\x04")
        with pytest.raises(ImportError, match="pip install.*memctl"):
            read_file_text(str(f))

    def test_xlsx_import_error(self, tmp_path):
        try:
            import openpyxl  # noqa: F401
            pytest.skip("openpyxl is installed — skip ImportError test")
        except ImportError:
            pass
        f = tmp_path / "test.xlsx"
        f.write_bytes(b"PK\x03\x04")
        with pytest.raises(ImportError, match="pip install.*memctl"):
            read_file_text(str(f))


# ---------------------------------------------------------------------------
# PDF extraction (requires pdftotext from poppler-utils)
# ---------------------------------------------------------------------------

HAS_PDFTOTEXT = shutil.which("pdftotext") is not None


@pytest.mark.skipif(not HAS_PDFTOTEXT, reason="pdftotext not installed")
class TestPDFExtraction:
    """Test PDF text extraction via pdftotext."""

    @pytest.fixture
    def sample_pdf(self, tmp_path):
        """Create a minimal PDF with known text content."""
        # Minimal valid PDF with embedded text
        pdf_content = (
            b"%PDF-1.0\n"
            b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
            b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
            b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R"
            b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
            b"4 0 obj<</Length 44>>\nstream\n"
            b"BT /F1 12 Tf 100 700 Td (Hello PDF) Tj ET\n"
            b"endstream\nendobj\n"
            b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
            b"xref\n0 6\n"
            b"0000000000 65535 f \n"
            b"0000000009 00000 n \n"
            b"0000000058 00000 n \n"
            b"0000000115 00000 n \n"
            b"0000000266 00000 n \n"
            b"0000000360 00000 n \n"
            b"trailer<</Size 6/Root 1 0 R>>\n"
            b"startxref\n431\n%%EOF\n"
        )
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(pdf_content)
        return str(pdf_path)

    def test_pdf_extraction(self, sample_pdf):
        text = read_file_text(sample_pdf)
        assert "Hello PDF" in text

    def test_pdf_empty_result_for_invalid(self, tmp_path):
        """Corrupt PDF should not crash, returns empty or partial text."""
        f = tmp_path / "corrupt.pdf"
        f.write_bytes(b"not a real pdf")
        # Should not raise, pdftotext returns non-zero
        text = read_file_text(str(f))
        assert isinstance(text, str)


# ---------------------------------------------------------------------------
# Office format extraction (requires pip install memctl[docs])
# ---------------------------------------------------------------------------

# These tests only run when the libraries are installed.
# To run: pip install memctl[docs] && pytest tests/test_extract.py -v

_HAS_DOCX = False
_HAS_PPTX = False
_HAS_OPENPYXL = False
_HAS_ODFPY = False

try:
    import docx  # noqa: F401
    _HAS_DOCX = True
except ImportError:
    pass
try:
    import pptx  # noqa: F401
    _HAS_PPTX = True
except ImportError:
    pass
try:
    import openpyxl  # noqa: F401
    _HAS_OPENPYXL = True
except ImportError:
    pass
try:
    import odf  # noqa: F401
    _HAS_ODFPY = True
except ImportError:
    pass


@pytest.mark.skipif(not _HAS_DOCX, reason="python-docx not installed")
class TestDocxExtraction:
    def test_extract_docx(self, tmp_path):
        from docx import Document
        doc = Document()
        doc.add_paragraph("First paragraph of the document.")
        doc.add_paragraph("Second paragraph with content.")
        path = tmp_path / "test.docx"
        doc.save(str(path))
        text = read_file_text(str(path))
        assert "First paragraph" in text
        assert "Second paragraph" in text

    def test_extract_docx_with_table(self, tmp_path):
        from docx import Document
        doc = Document()
        doc.add_paragraph("Header text")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"
        path = tmp_path / "table.docx"
        doc.save(str(path))
        text = read_file_text(str(path))
        assert "Header text" in text
        assert "A1" in text
        assert "|" in text  # table cells joined with |


@pytest.mark.skipif(not _HAS_PPTX, reason="python-pptx not installed")
class TestPptxExtraction:
    def test_extract_pptx(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Slide body content here"
        path = tmp_path / "test.pptx"
        prs.save(str(path))
        text = read_file_text(str(path))
        assert "[Slide 1]" in text
        assert "Slide Title" in text
        assert "Slide body content" in text


@pytest.mark.skipif(not _HAS_OPENPYXL, reason="openpyxl not installed")
class TestXlsxExtraction:
    def test_extract_xlsx(self, tmp_path):
        from openpyxl import Workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Age", "City"])
        ws.append(["Alice", 30, "Paris"])
        ws.append(["Bob", 25, "London"])
        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        text = read_file_text(str(path))
        assert "[Sheet: Data]" in text
        assert "Alice" in text
        assert "Paris" in text
        assert "|" in text  # cells joined with |


@pytest.mark.skipif(not _HAS_ODFPY, reason="odfpy not installed")
class TestOdfExtraction:
    def test_extract_odt(self, tmp_path):
        from odf.opendocument import OpenDocumentText
        from odf.text import P as OdfP
        doc = OpenDocumentText()
        p1 = OdfP()
        p1.addText("First ODF paragraph.")
        doc.text.addElement(p1)
        p2 = OdfP()
        p2.addText("Second ODF paragraph.")
        doc.text.addElement(p2)
        path = tmp_path / "test.odt"
        doc.save(str(path))
        text = read_file_text(str(path))
        assert "First ODF paragraph" in text
        assert "Second ODF paragraph" in text


# ---------------------------------------------------------------------------
# Integration: ingest.py uses extract.py
# ---------------------------------------------------------------------------

class TestIngestIntegration:
    """Verify ingest.py imports and uses extract.py correctly."""

    def test_ingestable_exts_includes_binary(self):
        """The ingest module's _INGESTABLE_EXTS must include binary formats."""
        from memctl.ingest import _INGESTABLE_EXTS
        assert ".docx" in _INGESTABLE_EXTS
        assert ".xlsx" in _INGESTABLE_EXTS
        assert ".pdf" in _INGESTABLE_EXTS
        assert ".pptx" in _INGESTABLE_EXTS

    def test_resolve_sources_finds_text_files(self, tmp_path):
        """resolve_sources should find text files in directories."""
        from memctl.ingest import resolve_sources
        (tmp_path / "readme.md").write_text("# Hi", encoding="utf-8")
        (tmp_path / "data.json").write_text("{}", encoding="utf-8")
        (tmp_path / "image.png").write_bytes(b"\x89PNG")  # not ingestable
        files = resolve_sources([str(tmp_path)])
        names = {os.path.basename(f) for f in files}
        assert "readme.md" in names
        assert "data.json" in names
        assert "image.png" not in names

    def test_resolve_sources_finds_binary_exts(self, tmp_path):
        """resolve_sources should include binary format files when walking dirs."""
        from memctl.ingest import resolve_sources
        (tmp_path / "doc.docx").write_bytes(b"PK\x03\x04")
        (tmp_path / "sheet.xlsx").write_bytes(b"PK\x03\x04")
        (tmp_path / "slides.pptx").write_bytes(b"PK\x03\x04")
        (tmp_path / "report.pdf").write_bytes(b"%PDF-1.0")
        files = resolve_sources([str(tmp_path)])
        names = {os.path.basename(f) for f in files}
        assert "doc.docx" in names
        assert "sheet.xlsx" in names
        assert "slides.pptx" in names
        assert "report.pdf" in names
