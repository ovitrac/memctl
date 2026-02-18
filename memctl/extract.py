"""
Text Extraction — Unified reader for text and binary document formats

Supports:
    Text files   .md .txt .rst .py .java .yaml .json .xml .html .csv  (direct read)
    Source code   .js .ts .go .rs .c .cpp .sh .sql .css …              (direct read)
    Office docs  .docx .odt                                           (python-docx / odfpy)
    Slides       .pptx .odp                                           (python-pptx / odfpy)
    Spreadsheets .xlsx .ods                                           (openpyxl / odfpy)
    PDF          .pdf                                                  (pdftotext via poppler)

Each extractor is optional: a missing library triggers a clear ImportError
with install instructions.  The public entry point is ``read_file_text(path)``.

Author: Olivier Vitrac, PhD, HDR | olivier.vitrac@adservio.fr | Adservio
"""

from __future__ import annotations

import logging
import subprocess
from pathlib import Path

logger = logging.getLogger(__name__)

# Extensions handled as binary (need extraction, not plain read)
BINARY_EXTS = frozenset({
    ".docx", ".odt",
    ".pptx", ".odp",
    ".xlsx", ".ods",
    ".pdf",
})

# All extensions considered ingestable (text + source code + binary)
ALL_INGESTABLE_EXTS = frozenset({
    # Text / markup
    ".md", ".txt", ".rst", ".csv", ".tsv",
    ".html", ".htm", ".xml", ".json", ".yaml", ".yml", ".toml",
    # Source code
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".go", ".rs", ".c", ".h", ".cpp", ".hpp",
    ".css", ".scss", ".less",
    ".sh", ".bash", ".zsh",
    ".sql", ".r", ".jl", ".lua", ".rb", ".php", ".swift", ".kt",
    ".dockerfile",
    # Binary (office / PDF)
    ".docx", ".odt",
    ".pptx", ".odp",
    ".xlsx", ".ods",
    ".pdf",
})


def read_file_text(path: str, *, encoding: str = "utf-8") -> str:
    """
    Read any supported file and return its text content.

    For text files: direct ``Path.read_text()``.
    For binary formats: extract text using the appropriate library.

    Args:
        path: File path (absolute or relative).
        encoding: Encoding for text files (default: utf-8).

    Returns:
        Extracted text content.

    Raises:
        ImportError: When a required extraction library is not installed.
        FileNotFoundError: When the file does not exist.
        ValueError: When the file extension is not supported.
    """
    p = Path(path)
    ext = p.suffix.lower()

    if ext in BINARY_EXTS:
        return _extract_binary(str(p), ext)

    # Text file: direct read
    return p.read_text(encoding=encoding, errors="replace")


# ---------------------------------------------------------------------------
# Binary format extractors
# ---------------------------------------------------------------------------


def _extract_binary(path: str, ext: str) -> str:
    """Dispatch to the appropriate binary extractor."""
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".odt":
        return _extract_odt(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".odp":
        return _extract_odp(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext == ".ods":
        return _extract_ods(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    raise ValueError(f"Unsupported binary format: {ext}")


# --- DOCX (python-docx) ---------------------------------------------------

def _extract_docx(path: str) -> str:
    """Extract text from a .docx file using python-docx."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required for .docx files. "
            "Install with: pip install python-docx   "
            "(or: pip install memctl[docs])"
        )
    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    # Also extract text from tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


# --- ODT (odfpy) ----------------------------------------------------------

def _extract_odt(path: str) -> str:
    """Extract text from an .odt file using odfpy."""
    try:
        from odf.opendocument import load as odf_load
        from odf.text import P as OdfP
        from odf import teletype
    except ImportError:
        raise ImportError(
            "odfpy is required for .odt files. "
            "Install with: pip install odfpy   "
            "(or: pip install memctl[docs])"
        )
    doc = odf_load(path)
    parts: list[str] = []
    for para in doc.getElementsByType(OdfP):
        text = teletype.extractText(para).strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


# --- PPTX (python-pptx) ---------------------------------------------------

def _extract_pptx(path: str) -> str:
    """Extract text from a .pptx file using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for .pptx files. "
            "Install with: pip install python-pptx   "
            "(or: pip install memctl[docs])"
        )
    prs = Presentation(path)
    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


# --- ODP (odfpy) ----------------------------------------------------------

def _extract_odp(path: str) -> str:
    """Extract text from an .odp file using odfpy."""
    try:
        from odf.opendocument import load as odf_load
        from odf.text import P as OdfP
        from odf import teletype
        from odf.draw import Page as OdfPage
    except ImportError:
        raise ImportError(
            "odfpy is required for .odp files. "
            "Install with: pip install odfpy   "
            "(or: pip install memctl[docs])"
        )
    doc = odf_load(path)
    parts: list[str] = []
    pages = doc.getElementsByType(OdfPage)
    for slide_num, page in enumerate(pages, 1):
        slide_texts: list[str] = []
        for para in page.getElementsByType(OdfP):
            text = teletype.extractText(para).strip()
            if text:
                slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


# --- XLSX (openpyxl) -------------------------------------------------------

def _extract_xlsx(path: str) -> str:
    """Extract text from a .xlsx file using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise ImportError(
            "openpyxl is required for .xlsx files. "
            "Install with: pip install openpyxl   "
            "(or: pip install memctl[docs])"
        )
    wb = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_text: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells)
            if line.strip(" |"):
                rows_text.append(line)
        if rows_text:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))
    wb.close()
    return "\n\n".join(parts)


# --- ODS (odfpy) -----------------------------------------------------------

def _extract_ods(path: str) -> str:
    """Extract text from an .ods file using odfpy."""
    try:
        from odf.opendocument import load as odf_load
        from odf.table import Table as OdfTable, TableRow, TableCell
        from odf import teletype
    except ImportError:
        raise ImportError(
            "odfpy is required for .ods files. "
            "Install with: pip install odfpy   "
            "(or: pip install memctl[docs])"
        )
    doc = odf_load(path)
    parts: list[str] = []
    for table in doc.getElementsByType(OdfTable):
        sheet_name = table.getAttribute("name") or "Sheet"
        rows_text: list[str] = []
        for row in table.getElementsByType(TableRow):
            cells: list[str] = []
            for cell in row.getElementsByType(TableCell):
                text = teletype.extractText(cell).strip()
                cells.append(text)
            line = " | ".join(cells)
            if line.strip(" |"):
                rows_text.append(line)
        if rows_text:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows_text))
    return "\n\n".join(parts)


# --- PDF (pdftotext via poppler) -------------------------------------------

def _extract_pdf(path: str) -> str:
    """Extract text from a PDF using pdftotext (poppler-utils).

    Requires the ``pdftotext`` binary from the ``poppler-utils`` system
    package (``sudo apt install poppler-utils`` on Debian/Ubuntu,
    ``brew install poppler`` on macOS).
    """
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", path, "-"],
            capture_output=True,
            text=True,
            timeout=60,
        )
    except FileNotFoundError:
        raise ImportError(
            "pdftotext (poppler-utils) is required for .pdf files. "
            "Install with: sudo apt install poppler-utils   "
            "(or: brew install poppler on macOS)"
        )
    except subprocess.TimeoutExpired:
        logger.warning("PDF extraction timed out for %s", path)
        return ""

    if result.returncode != 0:
        logger.warning(
            "pdftotext returned %d for %s: %s",
            result.returncode, path, result.stderr[:200],
        )
        return ""

    return result.stdout
